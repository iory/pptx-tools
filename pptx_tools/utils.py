import logging
from pathlib import Path

import coloredlogs
from eos import makedirs
import langdetect
from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from tqdm import tqdm

from pptx_tools.audio_utils import get_wave_duration
from pptx_tools.data import get_transparent_img_path
from pptx_tools.slide_info import get_slide_infos
from pptx_tools.slide_info import slide_info_to_dict
from pptx_tools.slide_transition import set_slide_duration
from pptx_tools.video_utils import add_video_duration
from pptx_tools.video_utils import max_video_duration


base_logger = logging.getLogger(__name__)


def xpath(el, query):
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    return etree.ElementBase.xpath(el, query, namespaces=nsmap)


def autoplay_media(media):
    el_id = xpath(media.element, './/p:cNvPr')[0].attrib['id']
    el_cnt = xpath(
        media.element.getparent().getparent().getparent(),
        './/p:timing//p:video//p:spTgt[@spid="%s"]' % el_id,
    )[0]
    cond = xpath(el_cnt.getparent().getparent(), './/p:cond')[0]
    cond.set('delay', '0')


def delete_slide_notes_text(slide_path, logger=None):
    if logger is None:
        logger = base_logger
    coloredlogs.install(level='DEBUG', logger=logger)

    presentation = Presentation(slide_path)

    for page, slide in tqdm(enumerate(presentation.slides, start=1),
                            total=len(presentation.slides)):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            slide.notes_slide.notes_text_frame.text = ''
    return presentation


def add_synthesize_audio(slide_path, outdir, logger=None,
                         voice_name=None,
                         tts='google',
                         slide_duration_offset=1.0):
    """Synthesizes speech from the pptx."""

    if logger is None:
        logger = base_logger
    coloredlogs.install(level='DEBUG', logger=logger)

    if tts == 'google':
        from pptx_tools.tts.google_tts import google_text_to_speech
        text_to_speech = google_text_to_speech
    elif tts == 'azure':
        from pptx_tools.tts.azure_tts import azure_text_to_speech
        text_to_speech = azure_text_to_speech
    else:
        raise ValueError("Not supported tts {}".format(tts))

    output_path = Path(outdir)
    makedirs(output_path)

    presentation = Presentation(slide_path)
    slide_infos = get_slide_infos(presentation)
    add_video_duration(slide_path, slide_infos)
    slide_info_dict = slide_info_to_dict(slide_infos)

    total_time = 0.0
    for page, slide in tqdm(enumerate(presentation.slides, start=1),
                            total=len(presentation.slides)):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            note_txt = slide.notes_slide.notes_text_frame.text
            note_txt = note_txt.replace('\n', ' ')
            wave_path = output_path / '{}.wav'.format(page)

            lang = langdetect.detect(note_txt)
            if lang not in ['en', 'ja']:
                logger.warning('Not supported language:{}'.format(lang))
                lang = 'en'
                logger.warning('Use english')
            if voice_name is None:
                if lang == 'en':
                    if tts == 'google':
                        voice_name = 'en-US-Wavenet-A'
                    elif tts == 'azure':
                        voice_name = 'en-US-BrandonNeural'
                    else:
                        raise RuntimeError('invalid tts')
                elif lang == 'ja':
                    if tts == 'google':
                        voice_name = 'ja-JP-Wavenet-C'
                    elif tts == 'azure':
                        voice_name = 'ja-JP-NanamiNeural'
                    else:
                        raise RuntimeError('invalid tts')
            text_to_speech(wave_path, note_txt,
                           voice_name=voice_name)
            duration = get_wave_duration(wave_path)
            video_duration = max_video_duration(slide_info_dict[page])
            duration = max(duration, video_duration)
            set_slide_duration(slide, duration + slide_duration_offset)
            total_time += duration + slide_duration_offset
            try:
                movie = slide.shapes.add_movie(
                    str(wave_path),
                    Inches(0), Inches(0), Inches(1.0), Inches(1.0),
                    poster_frame_image=str(get_transparent_img_path()),
                    mime_type='video/unknown')
            except AttributeError:
                msg = '[Synthesize Slide]: Skip slide number:{}'.format(page)
                logger.warning(msg)
                continue
            autoplay_media(movie)
    print('Total sound time: {}'.format(total_time))
    output_slide_path = output_path / Path(slide_path).name
    presentation.save(output_slide_path)
    return output_slide_path
