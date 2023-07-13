from collections import defaultdict

from pptx.enum.shapes import MSO_SHAPE_TYPE


def get_slide_infos(prs):
    slide_infos = []
    num_media = 0

    for num_slide, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                num_media += 1
                slide_infos.append({
                    'shape_id': shape.shape_id,
                    'filename': f"ppt/media/media{num_media}",
                    'num_slide': num_slide,
                    'num_media': num_media,
                })

    return slide_infos


def slide_info_to_dict(slide_infos):
    slide_info_dict = defaultdict(list)
    for info in slide_infos:
        slide_info_dict[info['num_slide']].append(info)
    return slide_info_dict


def get_slide_info(filename, infos):
    for info in infos:
        if filename in info['filename']:
            return info
