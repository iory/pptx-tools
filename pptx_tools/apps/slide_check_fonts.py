#!/usr/bin/env python

import argparse

from pptx import Presentation


def main():
    parser = argparse.ArgumentParser(description='Check fonts for powerpoint.')
    parser.add_argument("input", help="The input pptx file.")
    args = parser.parse_args()
    prs = Presentation(args.input)

    for num_slide, slide in enumerate(prs.slides, 1):
        fonts = set()
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.text is None:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name is None:
                        continue
                    fonts.add(run.font.name)
        if len(fonts) > 0:
            print('Slide {}: {}'.format(num_slide, fonts))


if __name__ == '__main__':
    main()
